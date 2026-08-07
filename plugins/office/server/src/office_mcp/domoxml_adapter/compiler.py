"""The only module that translates Office domain state to/from domOXML."""

import asyncio
import base64
import io
import mimetypes
import re
from dataclasses import dataclass
from typing import cast

import tinycss2
from bs4 import BeautifulSoup, Tag
from domoxml import Presentation, Slide, pptx_to_html
from domoxml.types import (
    CustomSize,
    Fonts,
    OutputFormat,
    Palette,
    Theme,
    Transition,
)
from domoxml.types import (
    SlideSize as DomSlideSize,
)
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from office_mcp.domain.html import parse_styles, sanitize_fragment, serialize_styles
from office_mcp.domain.state import PresentationSnapshot, StoredSlide
from office_mcp.errors import ErrorCode, OfficeError
from office_mcp.models.common import (
    CustomSlideSize,
    PresetSlideSize,
    SlideSize,
    SlideSizePreset,
    SlideTransition,
)


@dataclass(frozen=True)
class ImportedPresentation:
    slides: list[StoredSlide]
    size: SlideSize
    warnings: list[dict[str, object]]
    preservation: list[dict[str, object]]
    coverage: list[dict[str, object]]


def _size(size: SlideSize):
    if isinstance(size, PresetSlideSize):
        return DomSlideSize(size.preset.value)
    assert isinstance(size, CustomSlideSize)
    return CustomSize(width_in=size.width_in, height_in=size.height_in)


def _dimensions(size: SlideSize) -> tuple[float, float]:
    if isinstance(size, CustomSlideSize):
        return round(size.width_in, 4), round(size.height_in, 4)
    return {
        SlideSizePreset.WIDE_16_9: (13.333, 7.5),
        SlideSizePreset.STANDARD_4_3: (10.0, 7.5),
        SlideSizePreset.WIDE_16_10: (10.0, 6.25),
    }[size.preset]


def _model_size(width_px: int, height_px: int) -> SlideSize:
    presets = {
        (1280, 720): SlideSizePreset.WIDE_16_9,
        (960, 720): SlideSizePreset.STANDARD_4_3,
        (960, 600): SlideSizePreset.WIDE_16_10,
    }
    preset = presets.get((width_px, height_px))
    if preset:
        return PresetSlideSize(preset=preset)
    return CustomSlideSize(width_in=round(width_px / 96, 4), height_in=round(height_px / 96, 4))


def _slide_number(value: str) -> int | None:
    match = re.search(r"(?:^|/)slide(\d+)\.xml", value, re.IGNORECASE)
    return int(match.group(1)) if match else None


def _inline_css(html: str, css: str) -> str:
    soup = BeautifulSoup(html, "html.parser")
    rules = tinycss2.parse_stylesheet(css, skip_comments=True, skip_whitespace=True)
    for rule in rules:
        if rule.type != "qualified-rule":
            continue
        selector = tinycss2.serialize(rule.prelude).strip()
        declarations = parse_styles(tinycss2.serialize(rule.content))
        try:
            matches = soup.select(selector)
        except Exception:
            continue
        for tag in matches:
            existing = parse_styles(str(tag.get("style", "")))
            declarations_for_tag = dict(declarations)
            declarations_for_tag.update(existing)
            tag["style"] = serialize_styles(declarations_for_tag)
    return str(soup)


def _with_transition(html: str, transition: SlideTransition | None) -> str:
    if transition is None:
        return html
    soup = BeautifulSoup(html, "html.parser")
    root = soup.find(True)
    if isinstance(root, Tag):
        root["data-transition"] = transition.value
    return str(soup)


class DomOXMLAdapter:
    def __init__(self, *, max_concurrency: int = 2, timeout_seconds: float = 120.0) -> None:
        self._render_slots = asyncio.Semaphore(max(1, max_concurrency))
        self._timeout_seconds = timeout_seconds

    def _presentation(self, snapshot: PresentationSnapshot) -> Presentation:
        palette = snapshot.theme.palette
        fonts = snapshot.theme.fonts
        deck = Presentation(
            size=_size(snapshot.size),
            theme=Theme(
                palette=Palette(
                    background=palette.background,
                    foreground=palette.foreground,
                    accent=palette.accent,
                    muted=palette.muted,
                ),
                fonts=Fonts(heading=fonts.heading, body=fonts.body),
            ),
        )
        for slide in snapshot.slides:
            deck.add(
                Slide(
                    # domOXML's public extractor reads typed transition intent
                    # from the slide root. Inject it here so the public
                    # Slide.transition value survives the current render path.
                    html=_with_transition(slide.html, slide.transition),
                    transition=Transition(slide.transition.value) if slide.transition else None,
                    size=_size(slide.size) if slide.size else None,
                )
            )
        return deck

    async def render(
        self,
        snapshot: PresentationSnapshot,
        formats: set[OutputFormat],
        indices: set[int] | None = None,
    ):
        try:
            async with self._render_slots, asyncio.timeout(self._timeout_seconds):
                return await self._presentation(snapshot).arender(formats, indices=indices)
        except TimeoutError as exc:
            raise OfficeError(ErrorCode.RENDER_FAILED, "presentation rendering timed out") from exc
        except OfficeError:
            raise
        except Exception as exc:
            raise OfficeError(
                ErrorCode.RENDER_FAILED, "domOXML could not render the presentation"
            ) from exc

    async def export_pptx(self, snapshot: PresentationSnapshot) -> tuple[bytes, object]:
        if snapshot.imported_pptx_b64 and not snapshot.content_changed_after_import:
            return base64.b64decode(snapshot.imported_pptx_b64), None
        if snapshot.imported_pptx_b64 and snapshot.imported_preservation:
            raise OfficeError(
                ErrorCode.EXPORT_FAILED,
                "this edited import contains OOXML preservation fragments that the current "
                "domOXML public API cannot safely reattach; export is blocked to prevent "
                "silent data loss",
            )
        if not snapshot.slides:
            buffer = io.BytesIO()
            empty = PptxPresentation()
            width_in, height_in = _dimensions(snapshot.size)
            empty.slide_width = Inches(width_in)
            empty.slide_height = Inches(height_in)
            empty.save(buffer)
            return buffer.getvalue(), None
        if self.mixed_size_indices(snapshot):
            raise OfficeError(
                ErrorCode.EXPORT_FAILED,
                "PowerPoint files require one slide size; clear per-slide size overrides "
                "or make every effective slide size match the presentation size",
            )
        result = await self.render(snapshot, {OutputFormat.PPTX})
        if result.pptx is None:
            raise OfficeError(ErrorCode.EXPORT_FAILED, "domOXML produced no PPTX output")
        return result.pptx, result

    @staticmethod
    def mixed_size_indices(snapshot: PresentationSnapshot) -> set[int]:
        deck_dimensions = _dimensions(snapshot.size)
        return {
            index
            for index, slide in enumerate(snapshot.slides)
            if slide.size is not None and _dimensions(slide.size) != deck_dimensions
        }

    async def preview_pngs(
        self, snapshot: PresentationSnapshot, indices: set[int]
    ) -> tuple[bytes, ...]:
        if not indices:
            return ()
        result = await self.render(snapshot, {OutputFormat.PNG}, indices=indices)
        return result.pngs

    async def validate(self, snapshot: PresentationSnapshot, indices: set[int] | None = None):
        if not snapshot.slides:
            return None
        return await self.render(snapshot, {OutputFormat.PPTX}, indices=indices)

    async def import_pptx(self, data: bytes) -> ImportedPresentation:
        try:
            async with self._render_slots, asyncio.timeout(self._timeout_seconds):
                imported = await asyncio.to_thread(pptx_to_html, data)
        except TimeoutError as exc:
            raise OfficeError(ErrorCode.IMPORT_FAILED, "PPTX import timed out") from exc
        except Exception as exc:
            raise OfficeError(ErrorCode.IMPORT_FAILED, "domOXML could not import the PPTX") from exc
        assets = {
            asset.path: (
                mimetypes.guess_type(asset.path)[0] or "application/octet-stream",
                base64.b64encode(asset.data).decode(),
            )
            for asset in imported.assets
        }
        deck_size: SlideSize = (
            _model_size(imported.slides[0].width_px, imported.slides[0].height_px)
            if imported.slides
            else PresetSlideSize(preset=SlideSizePreset.WIDE_16_9)
        )
        slides: list[StoredSlide] = []
        for index, source in enumerate(imported.slides, start=1):
            # domOXML emits package assets as paths relative to its generated
            # stylesheet (for example ``../assets/image.png``). Resolve those
            # trusted, extracted bytes before parsing styles so the general
            # HTML policy never has to permit filesystem-like CSS URLs.
            source_html = source.html
            source_css = imported.css
            for path, (mime_type, encoded) in assets.items():
                data_uri = f"data:{mime_type};base64,{encoded}"
                normalized_path = path.lstrip("./")
                references = {
                    path,
                    normalized_path,
                    f"./{normalized_path}",
                    f"../{normalized_path}",
                }
                for reference in sorted(references, key=len, reverse=True):
                    source_html = source_html.replace(reference, data_uri)
                    source_css = source_css.replace(reference, data_uri)
            html = _inline_css(source_html, source_css)
            source_soup = BeautifulSoup(html, "html.parser")
            source_root = source_soup.find(True)
            transition_value = (
                str(source_root.get("data-transition"))
                if isinstance(source_root, Tag) and source_root.get("data-transition")
                else None
            )
            normalized, _ = sanitize_fragment(html, _preserve_domoxml_metadata=True)
            soup = BeautifulSoup(normalized, "html.parser")
            title = soup.find(["h1", "h2", "h3", "title"])
            name = (
                " ".join(title.stripped_strings)[:80]
                if isinstance(title, Tag)
                else f"Slide {index}"
            )
            slide_size = _model_size(source.width_px, source.height_px)
            slides.append(
                StoredSlide(
                    slide_id="",
                    name=name or f"Slide {index}",
                    html=normalized,
                    transition=SlideTransition(transition_value)
                    if transition_value in {item.value for item in SlideTransition}
                    else None,
                    size=slide_size if slide_size != deck_size else None,
                )
            )
        warnings: list[dict[str, object]] = [
            {
                "code": "DOMOXML_WARNING",
                "message": item.message,
                "slide_number": _slide_number(item.element),
                "element": item.element or None,
            }
            for item in imported.warnings
        ]
        preservation = [item.model_dump(mode="json") for item in imported.preserved]
        coverage: list[dict[str, object]] = [
            cast(
                dict[str, object],
                {
                    **item.model_dump(mode="json"),
                    "slide_number": _slide_number(item.element),
                },
            )
            for item in imported.coverage.items
        ]
        return ImportedPresentation(slides, deck_size, warnings, preservation, coverage)
