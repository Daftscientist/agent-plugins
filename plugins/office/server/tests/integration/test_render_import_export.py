import base64
import io
from typing import Any, cast

import pytest
from bs4 import BeautifulSoup, Tag
from mcp.server import MCPServer
from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from office_mcp.app import OfficeRuntime
from office_mcp.errors import ErrorCode, OfficeError
from office_mcp.inputs.pptx import validate_pptx
from office_mcp.models.common import PresetSlideSize, SlideSizePreset, SourceRetention
from office_mcp.models.element import ElementById, ElementMutation, ElementUpdateArgs
from office_mcp.models.presentation import (
    NewSlide,
    PresentationCreateArgs,
    PresentationExportArgs,
    PresentationOpenArgs,
    PresentationSource,
)
from office_mcp.models.preview import PresentationPreviewArgs, PreviewSlides
from office_mcp.models.slide import SlideInspectArgs, SlideInspectDetail
from office_mcp.models.validation import PresentationValidateArgs, ValidationDetail
from office_mcp.storage.protocols import LOCAL_SCOPE


async def test_real_domoxml_preview_validation_export_and_import(
    office: tuple[MCPServer[Any], OfficeRuntime],
) -> None:
    _, runtime = office
    created = await runtime.service.create(
        LOCAL_SCOPE,
        PresentationCreateArgs(
            name="Rendered",
            slides=[
                NewSlide(
                    name="Cover",
                    html=(
                        '<section style="width:100%;height:100%;background:#fff;padding:64px">'
                        '<h1 data-office-name="title" style="font-size:42px;color:#111">'
                        "Real render</h1></section>"
                    ),
                )
            ],
        ),
    )
    images, preview = await runtime.service.preview(
        LOCAL_SCOPE,
        PresentationPreviewArgs(
            presentation_id=created.presentation_id,
            selection=PreviewSlides(slide_ids=[created.slides[0].slide_id]),
        ),
    )
    assert preview.images[0].slide_ids == [created.slides[0].slide_id]
    image = Image.open(io.BytesIO(images[0]))
    assert image.width > image.height

    validation = await runtime.service.validate(
        LOCAL_SCOPE,
        PresentationValidateArgs(
            presentation_id=created.presentation_id, detail=ValidationDetail.FULL
        ),
    )
    assert validation.valid
    assert validation.coverage is not None

    exported = await runtime.service.export(
        LOCAL_SCOPE, PresentationExportArgs(presentation_id=created.presentation_id)
    )
    data = await runtime.output.read(LOCAL_SCOPE, exported.resource_uri)
    assert data is not None
    validate_pptx(data, runtime.config)
    assert exported.sha256 and exported.size_bytes == len(data)

    source = (
        "data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,"
        + base64.b64encode(data).decode()
    )
    imported = await runtime.service.open(
        LOCAL_SCOPE,
        PresentationOpenArgs(source=PresentationSource(uri=source, filename_hint="roundtrip.pptx")),
    )
    untouched = await runtime.service.export(
        LOCAL_SCOPE, PresentationExportArgs(presentation_id=imported.presentation_id)
    )
    untouched_bytes = await runtime.output.read(LOCAL_SCOPE, untouched.resource_uri)
    assert untouched_bytes == data


async def test_imported_dimensions_and_preservation_debt_remain_truthful(
    office: tuple[MCPServer[Any], OfficeRuntime],
) -> None:
    _, runtime = office
    source = PptxPresentation()
    slide = source.slides.add_slide(source.slide_layouts[6])
    chart_data = ChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("Series", (1, 2))  # pyright: ignore[reportUnknownMemberType]
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(4),
        chart_data,
    )
    slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(4), Inches(0.5)).text = "Before"
    buffer = io.BytesIO()
    source.save(buffer)
    encoded = base64.b64encode(buffer.getvalue()).decode()
    opened = await runtime.service.open(
        LOCAL_SCOPE,
        PresentationOpenArgs(
            source=PresentationSource(
                uri=(
                    "data:application/vnd.openxmlformats-officedocument."
                    f"presentationml.presentation;base64,{encoded}"
                )
            )
        ),
    )
    snapshot = await runtime.store.get(LOCAL_SCOPE, opened.presentation_id)
    assert snapshot.size == PresetSlideSize(preset=SlideSizePreset.STANDARD_4_3)
    assert snapshot.imported_preservation
    assert any(warning.code == "PRESERVED_SOURCE" for warning in opened.warnings)
    inspected = await runtime.service.slide_inspect(
        LOCAL_SCOPE,
        SlideInspectArgs(
            presentation_id=opened.presentation_id,
            slide_id=opened.slides[0].slide_id,
            detail=SlideInspectDetail.SOURCE,
        ),
    )
    assert inspected.html and "data-domoxml-" not in inspected.html
    structure = await runtime.service.slide_inspect(
        LOCAL_SCOPE,
        SlideInspectArgs(
            presentation_id=opened.presentation_id,
            slide_id=opened.slides[0].slide_id,
            detail=SlideInspectDetail.STRUCTURE,
        ),
    )
    text_node = next(item for item in structure.structure or [] if item.text == "Before")

    changed = await runtime.service.element_update(
        LOCAL_SCOPE,
        ElementUpdateArgs(
            presentation_id=opened.presentation_id,
            slide_id=opened.slides[0].slide_id,
            expected_revision=opened.revision,
            elements=[
                ElementMutation(element=ElementById(element_id=text_node.element_id), text="After")
            ],
        ),
    )
    changed_snapshot = await runtime.store.get(LOCAL_SCOPE, opened.presentation_id)
    changed_soup = BeautifulSoup(changed_snapshot.slides[0].html, "html.parser")
    changed_tag = changed_soup.find(attrs=cast(Any, {"data-office-id": text_node.element_id}))
    assert isinstance(changed_tag, Tag)
    assert all(
        "data-domoxml-text-payload" not in ancestor.attrs
        for ancestor in [changed_tag, *changed_tag.parents]
        if isinstance(ancestor, Tag)
    )
    validation = await runtime.service.validate(
        LOCAL_SCOPE,
        PresentationValidateArgs(
            presentation_id=opened.presentation_id,
            revision=changed.revision,
            detail=ValidationDetail.FULL,
        ),
    )
    assert any(warning.code == "PRESERVATION_DEBT" for warning in validation.warnings)
    assert validation.coverage and any(
        item.source_retention is SourceRetention.LOST and item.representation.value == "failed"
        for item in validation.coverage
    )
    with pytest.raises(OfficeError) as export_error:
        await runtime.service.export(
            LOCAL_SCOPE,
            PresentationExportArgs(
                presentation_id=opened.presentation_id, revision=changed.revision
            ),
        )
    assert export_error.value.code is ErrorCode.EXPORT_FAILED
    assert "silent data loss" in export_error.value.message


async def test_empty_deck_export_preserves_configured_slide_size(
    office: tuple[MCPServer[Any], OfficeRuntime],
) -> None:
    _, runtime = office
    created = await runtime.service.create(
        LOCAL_SCOPE,
        PresentationCreateArgs(
            name="Empty widescreen",
            size=PresetSlideSize(preset=SlideSizePreset.WIDE_16_9),
        ),
    )
    exported = await runtime.service.export(
        LOCAL_SCOPE, PresentationExportArgs(presentation_id=created.presentation_id)
    )
    pptx = await runtime.output.read(LOCAL_SCOPE, exported.resource_uri)
    assert pptx is not None
    package = PptxPresentation(io.BytesIO(pptx))
    package_dimensions = cast(Any, package)
    slide_width = int(package_dimensions.slide_width)
    slide_height = int(package_dimensions.slide_height)
    assert round(slide_width / 914_400, 3) == 13.333
    assert round(slide_height / 914_400, 3) == 7.5
